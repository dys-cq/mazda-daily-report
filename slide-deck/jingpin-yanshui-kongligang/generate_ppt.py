# -*- coding: utf-8 -*-
"""
Generate a professional PowerPoint presentation for Job Competition Presentation
Using the content from the prompts and outline.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
prs.slide_height = Inches(7.5)

# Color scheme
NAVY_BLUE = RGBColor(26, 54, 93)    # #1a365d
WHITE = RGBColor(255, 255, 255)      # #ffffff
GOLD = RGBColor(214, 158, 46)        # #d69e2e
LIGHT_GRAY = RGBColor(247, 250, 252) # #f7fafc

def set_background(slide, color=NAVY_BLUE):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background(slide, NAVY_BLUE)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.font.name = 'Microsoft YaHei'
    
    # Add subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    sub_frame = sub_box.text_frame
    sub_para = sub_frame.paragraphs[0]
    sub_para.text = subtitle
    sub_para.alignment = PP_ALIGN.CENTER
    sub_para.font.size = Pt(32)
    sub_para.font.color.rgb = GOLD
    sub_para.font.name = 'Microsoft YaHei'
    
    return slide

def add_content_slide(prs, title, content_items, subtitle=None):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background(slide, WHITE)
    
    # Add title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY_BLUE
    title_bar.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.font.name = 'Microsoft YaHei'
    
    # Add subtitle if provided
    y_start = 1.5
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.5))
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = GOLD
        sub_para.font.name = 'Microsoft YaHei'
        y_start = 2.0
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_start), Inches(12), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()
        para.text = chr(8226) + " " + item  # Unicode bullet point
        para.font.size = Pt(24)
        para.font.color.rgb = NAVY_BLUE
        para.font.name = 'Microsoft YaHei'
        para.space_after = Pt(14)
    
    return slide

def add_three_pillar_slide(prs, title, pillars):
    """Add a three-pillar slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    
    # Add title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY_BLUE
    title_bar.line.fill.background()
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.font.name = 'Microsoft YaHei'
    
    # Three pillars
    pillar_width = 3.8
    gap = 0.3
    start_x = 0.7
    
    for i, (pillar_title, pillar_items) in enumerate(pillars):
        x = start_x + i * (pillar_width + gap)
        
        # Pillar background
        pillar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), pillar_width, Inches(5))
        pillar.fill.solid()
        pillar.fill.fore_color.rgb = LIGHT_GRAY
        pillar.line.color.rgb = GOLD
        
        # Pillar title
        title_box = slide.shapes.add_textbox(x + 0.2, Inches(1.7), pillar_width - 0.4, Inches(0.8))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = str(i+1) + ". " + pillar_title
        title_para.font.size = Pt(22)
        title_para.font.bold = True
        title_para.font.color.rgb = NAVY_BLUE
        title_para.font.name = 'Microsoft YaHei'
        
        # Pillar content
        content_box = slide.shapes.add_textbox(x + 0.2, Inches(2.5), pillar_width - 0.4, Inches(3.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for j, item in enumerate(pillar_items):
            if j == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()
            para.text = chr(8226) + " " + item
            para.font.size = Pt(18)
            para.font.color.rgb = NAVY_BLUE
            para.font.name = 'Microsoft YaHei'
            para.space_after = Pt(8)
    
    return slide

def add_closing_slide(prs, main_text, subtitle, name):
    """Add a closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY_BLUE)
    
    # Add main text
    main_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    main_frame = main_box.text_frame
    main_para = main_frame.paragraphs[0]
    main_para.text = main_text
    main_para.alignment = PP_ALIGN.CENTER
    main_para.font.size = Pt(54)
    main_para.font.bold = True
    main_para.font.color.rgb = WHITE
    main_para.font.name = 'Microsoft YaHei'
    
    # Add subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(1))
    sub_frame = sub_box.text_frame
    sub_para = sub_frame.paragraphs[0]
    sub_para.text = subtitle
    sub_para.alignment = PP_ALIGN.CENTER
    sub_para.font.size = Pt(28)
    sub_para.font.color.rgb = GOLD
    sub_para.font.name = 'Microsoft YaHei'
    
    # Add name
    name_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(1))
    name_frame = name_box.text_frame
    name_para = name_frame.paragraphs[0]
    name_para.text = name
    name_para.alignment = PP_ALIGN.CENTER
    name_para.font.size = Pt(32)
    name_para.font.bold = True
    name_para.font.color.rgb = WHITE
    name_para.font.name = 'Microsoft YaHei'
    
    return slide

def add_vision_slide(prs, title, levels):
    """Add a vision slide with three levels"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY_BLUE
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.font.name = 'Microsoft YaHei'
    
    # Three levels
    for i, (level_title, level_items) in enumerate(levels):
        y = 1.8 + i * 2
        level_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.333), Inches(1.7))
        level_frame = level_box.text_frame
        level_frame.word_wrap = True
        
        # Level title
        para = level_frame.paragraphs[0]
        para.text = str(i+1) + ". " + level_title
        para.font.size = Pt(24)
        para.font.bold = True
        para.font.color.rgb = GOLD
        para.font.name = 'Microsoft YaHei'
        
        # Level items
        for j, item in enumerate(level_items):
            p = level_frame.add_paragraph()
            p.text = "   " + chr(8226) + " " + item
            p.font.size = Pt(20)
            p.font.color.rgb = NAVY_BLUE
            p.font.name = 'Microsoft YaHei'
            p.space_after = Pt(6)
    
    return slide

def add_achievements_slide(prs, title, subtitle, sections):
    """Add an achievements slide with three sections"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY_BLUE
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.font.name = 'Microsoft YaHei'
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.5))
    sub_frame = sub_box.text_frame
    sub_para = sub_frame.paragraphs[0]
    sub_para.text = subtitle
    sub_para.font.size = Pt(24)
    sub_para.font.color.rgb = GOLD
    sub_para.font.name = 'Microsoft YaHei'
    
    # Three sections
    for i, (sec_title, sec_items) in enumerate(sections):
        x = 0.5 + i * 4.3
        sec_box = slide.shapes.add_textbox(x, Inches(2.2), Inches(4), Inches(4.5))
        sec_frame = sec_box.text_frame
        sec_frame.word_wrap = True
        
        # Section title
        para = sec_frame.paragraphs[0]
        para.text = sec_title
        para.font.size = Pt(24)
        para.font.bold = True
        para.font.color.rgb = NAVY_BLUE
        para.font.name = 'Microsoft YaHei'
        
        # Section items
        for j, item in enumerate(sec_items):
            p = sec_frame.add_paragraph()
            p.text = "  - " + item
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.font.name = 'Microsoft YaHei'
            p.space_after = Pt(8)
    
    return slide

# Generate slides
print("Generating presentation...")

# Slide 1: Cover
add_title_slide(prs, "竞聘演讲", "演讲人：孔立刚")
print("[OK] Slide 1: Cover")

# Slide 2: Personal Introduction
add_content_slide(prs, "个人介绍", [
    "孔立刚",
    "21 年从业经验",
    "资深汽车行业售后服务专家",
    "专注于售后管理体系建设与优化",
    "擅长数据驱动的业务决策与团队管理"
])
print("[OK] Slide 2: Personal Introduction")

# Slide 3: Strengths
add_content_slide(prs, "个人优势分析 (Strengths)", [
    "丰富的行业经验：21 年深耕汽车售后领域",
    "系统化思维：善于构建可视化管理体系",
    "技术创新能力：熟练运用自动化工具提升效率",
    "实战业绩突出：多个成功案例验证能力"
])
print("[OK] Slide 3: Strengths")

# Slide 4: Weaknesses
add_content_slide(prs, "个人劣势分析 (Weaknesses)", [
    "新技术学习速度需加快",
    "跨部门协调经验有待加强",
    "高层战略视野需要进一步拓展"
])
print("[OK] Slide 4: Weaknesses")

# Slide 5: Achievements - Case Study
sections = [
    ("项目背景", ["传统售后管理数据分散", "决策滞后"]),
    ("解决方案", ["搭建可视化数据管理平台", "实现关键指标实时监控", "建立自动化报表系统"]),
    ("取得成效", ["管理效率提升 40%", "客户满意度显著提高", "团队协作更加顺畅"])
]
add_achievements_slide(prs, "主要业绩成果", "重庆金团售后可视化管理成功案例", sections)
print("[OK] Slide 5: Achievements")

# Slide 6: Automation Tools
add_content_slide(prs, "业绩成果 - 自动化工具应用", [
    "数据分析自动化：减少人工统计时间 70%",
    "报表生成自动化：日报、周报自动生成",
    "预警机制自动化：异常数据实时提醒",
    "流程优化自动化：标准化作业流程"
])
print("[OK] Slide 6: Automation Tools")

# Slide 7: Competition Reasons
add_content_slide(prs, "竞聘理由", [
    "经验匹配：21 年经验完全契合岗位需求",
    "业绩支撑：成功案例证明能力",
    "愿景一致：与公司发展目标高度认同",
    "责任担当：愿意承担更大责任，创造更大价值"
])
print("[OK] Slide 7: Competition Reasons")

# Slide 8: Core Strategy
pillars = [
    ("数据驱动", ["用数据说话", "科学决策"]),
    ("客户至上", ["以提升客户体验为核心"]),
    ("团队赋能", ["激发团队潜能", "共同成长"])
]
add_three_pillar_slide(prs, "核心工作思路", pillars)
print("[OK] Slide 8: Core Strategy")

# Slide 9: Phase 1
add_content_slide(prs, "业务规划 - 第一阶段", [
    "夯实基础 (1-3 个月)",
    "梳理现有业务流程",
    "建立标准化作业规范",
    "完善数据采集体系",
    "搭建基础可视化平台"
])
print("[OK] Slide 9: Phase 1")

# Slide 10: Phase 2
add_content_slide(prs, "业务规划 - 第二阶段", [
    "优化提升 (4-9 个月)",
    "深化数据分析应用",
    "优化关键业务指标",
    "推广自动化管理工具",
    "建立持续改进机制"
])
print("[OK] Slide 10: Phase 2")

# Slide 11: Phase 3
add_content_slide(prs, "业务规划 - 第三阶段", [
    "创新突破 (10-18 个月)",
    "探索创新业务模式",
    "打造行业标杆案例",
    "形成可复制管理经验",
    "实现业务跨越式发展"
])
print("[OK] Slide 11: Phase 3")

# Slide 12-15: Team Management (4 slides)
team_slides = [
    ("要点 1 - 建立清晰的目标体系", [
        "设定明确的团队目标",
        "分解个人绩效指标",
        "定期回顾与调整",
        "确保目标与战略对齐"
    ]),
    ("要点 2 - 打造学习型组织", [
        "建立定期培训机制",
        "鼓励知识分享与交流",
        "支持成员技能提升",
        "营造持续学习氛围"
    ]),
    ("要点 3 - 强化沟通与协作", [
        "建立高效沟通渠道",
        "促进跨部门协作",
        "定期团队建设活动",
        "打造开放包容文化"
    ]),
    ("要点 4 - 激励机制与关怀", [
        "建立公平激励机制",
        "认可与奖励优秀表现",
        "关注成员职业发展",
        "营造温暖团队氛围"
    ])
]

for i, (subtitle, items) in enumerate(team_slides):
    add_content_slide(prs, "团队管理设想", items, subtitle)
    print("[OK] Slide " + str(12+i) + ": Team Management " + str(i+1))

# Slide 16: Vision
levels = [
    ("公司层面", [
        "助力公司战略目标实现",
        "推动公司数字化转型",
        "提升公司市场竞争力",
        "为公司创造持续价值"
    ]),
    ("管理层面", [
        "建立行业领先的售后管理标准",
        "形成可复制的管理方法论",
        "实现管理效率最大化",
        "成为行业管理标杆"
    ]),
    ("团队层面", [
        "培养高素质专业人才",
        "打造高绩效战斗团队",
        "实现团队成员共同成长",
        "成为行业人才摇篮"
    ])
]
add_vision_slide(prs, "愿景", levels)
print("[OK] Slide 16: Vision")

# Slide 17: Closing
add_closing_slide(prs, "感谢聆听", "期待与您共创美好未来", "孔立刚")
print("[OK] Slide 17: Closing")

# Save presentation
output_path = "C:/Users/Administrator/.openclaw/workspace/slide-deck/jingpin-yanshui-kongligang/竞聘演讲 - 孔立刚.pptx"
prs.save(output_path)
print("")
print("=" * 60)
print("Presentation saved to: " + output_path)
print("Total slides: 17")
print("=" * 60)
