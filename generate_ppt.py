#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
马自达创驰蓝天技术 PPT 生成器
生成一份适合普通人理解的演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_ppt():
    # 创建演示文稿
    prs = Presentation()
    
    # 设置 16:9 布局
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 定义颜色
    SKY_BLUE = RGBColor(0, 122, 204)      # 蓝天蓝
    DARK_BLUE = RGBColor(0, 51, 102)      # 深蓝
    WHITE = RGBColor(255, 255, 255)       # 白色
    LIGHT_GRAY = RGBColor(240, 240, 240)  # 浅灰
    TEXT_DARK = RGBColor(50, 50, 50)      # 深色文字
    ACCENT_ORANGE = RGBColor(255, 102, 0) # 强调橙色
    
    def add_title_slide(prs, title, subtitle):
        """添加封面页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        
        # 背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = SKY_BLUE
        bg.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(32)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(prs, title, content_items, notes=None):
        """添加内容页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 顶部标题栏
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = SKY_BLUE
        header.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(11), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # 内容
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            if isinstance(item, dict):
                # 带层级的内容
                p.text = item.get('text', '')
                p.font.size = Pt(item.get('size', 24))
                p.font.color.rgb = TEXT_DARK
                if item.get('bold'):
                    p.font.bold = True
                if item.get('indent'):
                    p.level = item['indent']
            else:
                p.text = item
                p.font.size = Pt(24)
                p.font.color.rgb = TEXT_DARK
                p.space_after = Pt(12)
        
        # 添加备注
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
        
        return slide
    
    def add_comparison_slide(prs, title, comparisons):
        """添加对比页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 顶部标题栏
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = SKY_BLUE
        header.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(11), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # 创建对比表格
        table_rows = len(comparisons) + 1
        table_cols = 4
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(12.333)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
        
        # 设置列宽
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(3.5)
        table.columns[2].width = Inches(3.5)
        table.columns[3].width = Inches(2.333)
        
        # 表头
        headers = ['项目', '传统技术', '创驰蓝天', '改善']
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.CENTER
        
        # 数据行
        for row_idx, comp in enumerate(comparisons, 1):
            for col_idx, value in enumerate(comp):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(16)
                    paragraph.alignment = PP_ALIGN.CENTER
                    if col_idx == 3:  # 改善列用橙色强调
                        paragraph.font.color.rgb = ACCENT_ORANGE
                        paragraph.font.bold = True
        
        return slide
    
    def add_summary_slide(prs, title, items):
        """添加总结页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 顶部标题栏
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = DARK_BLUE
        header.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(11), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # 内容 - 使用图标式布局
        emojis = ['🔧', '🔧', '⚙️', '🛡️', '🏎️', '💚', '💰', '🎯']
        y_pos = 2.0
        
        for i, item in enumerate(items):
            emoji = emojis[i] if i < len(emojis) else '•'
            
            # 添加圆形背景
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.8), Inches(y_pos),
                Inches(0.8), Inches(0.8)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = SKY_BLUE
            circle.line.fill.background()
            
            # 添加 emoji
            emoji_box = slide.shapes.add_textbox(
                Inches(0.9), Inches(y_pos + 0.15),
                Inches(0.6), Inches(0.5)
            )
            tf = emoji_box.text_frame
            p = tf.paragraphs[0]
            p.text = emoji
            p.font.size = Pt(28)
            p.alignment = PP_ALIGN.CENTER
            
            # 添加文字
            text_box = slide.shapes.add_textbox(
                Inches(1.8), Inches(y_pos + 0.2),
                Inches(11), Inches(0.5)
            )
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(22)
            p.font.color.rgb = TEXT_DARK
            
            y_pos += 0.9
        
        return slide
    
    def add_end_slide(prs, title, subtitle):
        """添加结束页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 背景
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BLUE
        bg.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(32)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    # ========== 开始创建幻灯片 ==========
    
    # 1. 封面
    add_title_slide(prs, "马自达创驰蓝天技术", "让驾驶更环保、更有趣")
    
    # 2. 什么是创驰蓝天技术
    add_content_slide(prs, "什么是创驰蓝天技术？", [
        "🔹 创驰蓝天（SKYACTIV）是马自达公司推出的全套汽车技术",
        "🔹 2010 年首次发布",
        "🔹 目标：既环保又有驾驶乐趣",
        "🔹 不是单一技术，而是包含发动机、变速箱、车身、底盘的完整系统",
        {"text": "简单说：让车更省油、更安全、更好开！", "bold": True, "size": 26}
    ])
    
    # 3. 为什么需要这项技术
    add_content_slide(prs, "为什么需要这项技术？", [
        {"text": "传统汽车的困境：", "bold": True, "size": 28},
        "  ❌ 油耗高 → 花钱多、污染大",
        "  ❌ 动力弱 → 开起来没意思",
        "  ❌ 车身重 → 不灵活、不安全",
        {"text": "创驰蓝天的解决方案：", "bold": True, "size": 28},
        "  ✅ 高效燃烧 → 省油",
        "  ✅ 轻量化 → 灵活又安全",
        "  ✅ 优化匹配 → 好开又舒适"
    ])
    
    # 4. 汽油发动机
    add_content_slide(prs, "核心技术：汽油发动机 SKYACTIV-G", [
        {"text": "超高压缩比 13:1", "bold": True, "size": 28},
        "  （普通车约 10:1）",
        "  → 汽油燃烧更充分",
        "  → 动力提升 15%，油耗降低 15%",
        {"text": "关键技术：", "bold": True, "size": 28},
        "  • 4-2-1 排气系统：减少废气残留",
        "  • 特殊活塞设计：让燃烧更均匀",
        "  • 直喷技术：燃油精准喷射",
        {"text": "比喻：像把柴火塞进炉子压得更紧，烧得更旺！", "size": 22}
    ])
    
    # 5. 柴油发动机
    add_content_slide(prs, "核心技术：柴油发动机 SKYACTIV-D", [
        {"text": "超低压缩比 14:1", "bold": True, "size": 28},
        "  （传统柴油车 16-18:1）",
        "  → 燃烧压力更低，噪音更小",
        "  → 无需昂贵的尾气处理系统",
        {"text": "优势：", "bold": True, "size": 28},
        "  ✅ 更安静",
        "  ✅ 更省油",
        "  ✅ 更环保",
        "  ✅ 维护成本低"
    ])
    
    # 6. 变速箱
    add_content_slide(prs, "核心技术：变速箱 SKYACTIV-DRIVE", [
        {"text": "融合三种变速箱的优点：", "bold": True, "size": 26},
        "  📌 AT 自动变速箱 → 平顺可靠",
        "  📌 CVT 无级变速 → 省油",
        "  📌 DCT 双离合 → 换挡快",
        {"text": "特点：", "bold": True, "size": 26},
        "  • 全速域锁止：减少动力损失",
        "  • 智能换挡：理解你的驾驶意图"
    ])
    
    # 7. 车身
    add_content_slide(prs, "核心技术：车身 SKYACTIV-BODY", [
        {"text": "轻量化设计：", "bold": True, "size": 28},
        "  • 使用高强度钢材",
        "  • 重量减轻 8%",
        "  • 刚性提升 30%",
        {"text": "安全性能：", "bold": True, "size": 28},
        "  • 多路径碰撞吸收：撞击时分散冲击力",
        "  • 环形结构：保护乘员舱",
        "  • 碰撞测试五星评级",
        {"text": "比喻：像鸡蛋壳，又轻又硬！", "size": 22}
    ])
    
    # 8. 底盘
    add_content_slide(prs, "核心技术：底盘 SKYACTIV-CHASSIS", [
        {"text": "操控优化：", "bold": True, "size": 28},
        "  • 悬架几何优化：转弯更稳",
        "  • 电动助力转向：手感自然",
        "  • 减震器调校：舒适与运动兼顾",
        {"text": "驾驶感受：", "bold": True, "size": 28},
        "  ✅ 指向精准",
        "  ✅ 过弯稳定",
        "  ✅ 长途不累",
        "  ✅ 城市灵活"
    ])
    
    # 9. 实际效果对比
    add_comparison_slide(prs, "实际效果对比", [
        ["马自达 3 油耗", "7.5L/100km", "6.0L/100km", "↓ 20%"],
        ["CX-5 油耗", "9.0L/100km", "7.0L/100km", "↓ 22%"],
        ["动力输出", "基准", "提升 15%", "↑ 15%"],
        ["车身重量", "基准", "减轻 8%", "↓ 8%"],
    ])
    
    # 10. 环保贡献
    add_content_slide(prs, "环保贡献", [
        {"text": "减少排放：", "bold": True, "size": 28},
        "  🌱 CO₂排放降低 15-20%",
        "  🌱 氮氧化物减少 30%",
        "  🌱 颗粒物排放更低",
        {"text": "能源效率：", "bold": True, "size": 28},
        "  • 燃油利用率提高",
        "  • 减少对石油依赖",
        "  • 为电动化过渡做准备"
    ])
    
    # 11. 搭载车型
    add_content_slide(prs, "搭载车型", [
        {"text": "轿车系列：", "bold": True, "size": 28},
        "  🚗 马自达 3（昂克赛拉）",
        "  🚗 马自达 6（阿特兹）",
        {"text": "SUV 系列：", "bold": True, "size": 28},
        "  🚙 CX-3、CX-4、CX-5",
        "  🚙 CX-8、CX-9",
        {"text": "跑车：", "bold": True, "size": 28},
        "  🏎️ MX-5（部分市场）"
    ])
    
    # 12. 品牌理念
    add_content_slide(prs, "品牌理念：「人马一体」", [
        {"text": "核心思想：", "bold": True, "size": 28},
        "  • 车是身体的延伸",
        "  • 驾驶应该是享受，不是负担",
        "  • 技术为人服务，不是人为技术服务",
        {"text": "品牌精神：", "bold": True, "size": 28},
        "  「Zoom-Zoom」",
        "  孩子看到车加速时的兴奋声",
        {"text": "马自达相信：好车应该让人开心！", "bold": True, "size": 24}
    ])
    
    # 13. 总结
    add_summary_slide(prs, "总结", [
        "SKYACTIV-G — 高效汽油发动机",
        "SKYACTIV-D — 清洁柴油发动机",
        "SKYACTIV-DRIVE — 智能变速箱",
        "SKYACTIV-BODY — 轻量安全车身",
        "SKYACTIV-CHASSIS — 精准操控底盘",
        "💚 环保：更低油耗、更少排放",
        "💰 经济：省油就是省钱",
        "🎯 乐趣：好开才是好车"
    ])
    
    # 14. 结束页
    add_end_slide(prs, "谢谢观看", "有问题吗？")
    
    # 保存文件
    output_path = "C:\\Users\\Administrator\\.openclaw\\workspace\\马自达创驰蓝天技术.pptx"
    prs.save(output_path)
    print(f"PPT 已生成：{output_path}")
    print(f"共 {len(prs.slides)} 页幻灯片")
    
    return output_path

if __name__ == "__main__":
    create_ppt()
